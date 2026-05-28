"""
export_service.py — Generates Excel, PDF, and PowerPoint reports.
All exports go to the EXPORTS_DIR and return a Path to the file.
"""

import copy
import io
import logging
from pathlib import Path
from datetime import datetime

import pandas as pd
import numpy as np
import plotly.graph_objects as go

from services.domain_service import DomainConfig
from services.kpi_service import KPIResult
from config import EXPORTS_DIR, CHART_COLORS
from utils.branding import BrandTheme, hex_to_rgb, summary_bullets

logger = logging.getLogger(__name__)


def _timestamp() -> str:
    return datetime.now().strftime("%Y%m%d_%H%M%S")


class ExportService:

    def _resolve_brand(self, brand: BrandTheme | None) -> BrandTheme:
        return brand or BrandTheme()

    def _styled_figure(self, fig: go.Figure, brand: BrandTheme) -> go.Figure:
        styled = copy.deepcopy(fig)
        palette = brand.palette or list(CHART_COLORS)
        styled.update_layout(
            colorway=palette,
            font=dict(color="#111827", family="Inter, Arial, sans-serif"),
            paper_bgcolor="#FFFFFF",
            plot_bgcolor="#FAFAFA",
        )
        return styled

    def _figure_png(self, fig: go.Figure, width: int, height: int) -> bytes:
        return fig.to_image(
            format="png", width=width, height=height, scale=2, engine="kaleido"
        )

    def export_csv(
        self,
        df: pd.DataFrame,
        domain: DomainConfig,
        filename: str | None = None,
    ) -> Path:
        """Export cleaned dataset as CSV."""
        filename = filename or f"cleaned_{domain.key}_{_timestamp()}.csv"
        path = EXPORTS_DIR / filename
        df.to_csv(path, index=False)
        logger.info("CSV exported: %s", path)
        return path

    # ── Excel ──────────────────────────────────────────────────────────────────
    def export_excel(
        self,
        df: pd.DataFrame,
        kpis: list[KPIResult],
        domain: DomainConfig,
        filename: str | None = None,
        brand: BrandTheme | None = None,
    ) -> Path:
        brand = self._resolve_brand(brand)
        filename = filename or f"analytics_{domain.key}_{_timestamp()}.xlsx"
        path = EXPORTS_DIR / filename

        with pd.ExcelWriter(path, engine="xlsxwriter") as writer:
            wb = writer.book

            header_fmt = wb.add_format({
                "bold": True, "bg_color": brand.primary, "font_color": "#FFFFFF",
                "border": 1, "align": "center",
            })
            kpi_label_fmt = wb.add_format({"bold": True, "font_color": "#374151"})
            kpi_value_fmt = wb.add_format({
                "num_format": "#,##0.00", "bold": True, "font_color": brand.primary,
            })
            num_fmt = wb.add_format({"num_format": "#,##0.00"})

            ws_kpi = wb.add_worksheet("KPI Summary")
            ws_kpi.set_column("A:A", 28)
            ws_kpi.set_column("B:B", 20)
            title = brand.company_name or f"{domain.label} — KPI Summary"
            ws_kpi.write(0, 0, title, wb.add_format({"bold": True, "font_size": 14}))
            ws_kpi.write(2, 0, "Metric", header_fmt)
            ws_kpi.write(2, 1, "Value", header_fmt)
            for i, kpi in enumerate(kpis):
                ws_kpi.write(3 + i, 0, kpi.name, kpi_label_fmt)
                ws_kpi.write(3 + i, 1, kpi.formatted, kpi_value_fmt)

            if brand.logo_bytes:
                logo_sheet = wb.add_worksheet("Brand")
                logo_sheet.set_column("A:A", 24)
                logo_sheet.insert_image("A1", "logo.png", {"image_data": io.BytesIO(brand.logo_bytes), "x_scale": 0.35, "y_scale": 0.35})
                logo_sheet.write("C1", brand.company_name or domain.label)
                logo_sheet.write("C2", f"Primary: {brand.primary}")
                logo_sheet.write("C3", f"Secondary: {brand.secondary}")
                logo_sheet.write("C4", f"Accent: {brand.accent}")

            df.to_excel(writer, sheet_name="Cleaned Data", index=False)
            ws_data = writer.sheets["Cleaned Data"]
            for col_num, col_name in enumerate(df.columns):
                ws_data.write(0, col_num, col_name, header_fmt)
                ws_data.set_column(col_num, col_num, 18)

            numeric_cols = df.select_dtypes(include=[np.number]).columns
            if len(numeric_cols):
                stats_df = df[numeric_cols].describe().T.round(3)
                stats_df.to_excel(writer, sheet_name="Statistics")

        logger.info("Excel exported: %s", path)
        return path

    # ── PDF ────────────────────────────────────────────────────────────────────
    def export_pdf(
        self,
        domain: DomainConfig,
        kpis: list[KPIResult],
        summary_text: str,
        charts: list[tuple[str, go.Figure]] | None = None,
        filename: str | None = None,
        brand: BrandTheme | None = None,
    ) -> Path:
        from reportlab.lib.pagesizes import A4
        from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
        from reportlab.lib.units import mm
        from reportlab.lib import colors
        from reportlab.platypus import (
            SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, HRFlowable, Image,
        )

        brand = self._resolve_brand(brand)
        filename = filename or f"report_{domain.key}_{_timestamp()}.pdf"
        path = EXPORTS_DIR / filename

        doc = SimpleDocTemplate(
            str(path), pagesize=A4,
            rightMargin=18 * mm, leftMargin=18 * mm,
            topMargin=18 * mm, bottomMargin=18 * mm,
        )
        styles = getSampleStyleSheet()
        BRAND = colors.HexColor(brand.primary)
        SECONDARY = colors.HexColor(brand.secondary)
        ACCENT = colors.HexColor(brand.accent)
        DARK = colors.HexColor("#111827")

        title_style = ParagraphStyle(
            "Title", parent=styles["Title"],
            textColor=BRAND, fontSize=22, spaceAfter=4,
        )
        subtitle_style = ParagraphStyle(
            "Sub", parent=styles["Normal"],
            textColor=DARK, fontSize=11, spaceAfter=12,
        )
        body_style = ParagraphStyle(
            "Body", parent=styles["Normal"],
            textColor=DARK, fontSize=10, leading=16,
        )
        kpi_label = ParagraphStyle(
            "KPILabel", parent=styles["Normal"],
            textColor=colors.HexColor("#6B7280"), fontSize=9,
        )
        kpi_value = ParagraphStyle(
            "KPIValue", parent=styles["Normal"],
            textColor=BRAND, fontSize=14, fontName="Helvetica-Bold",
        )

        story = []

        if brand.logo_bytes:
            logo_stream = io.BytesIO(brand.logo_bytes)
            story.append(Image(logo_stream, width=28 * mm, height=14 * mm))
            story.append(Spacer(1, 6))

        report_title = brand.company_name or f"{domain.label} Report"
        story.append(Paragraph(report_title, title_style))
        story.append(Paragraph(
            f"Generated {datetime.now().strftime('%B %d, %Y')} | {domain.label}",
            subtitle_style,
        ))
        story.append(HRFlowable(width="100%", thickness=2, color=BRAND, spaceAfter=12))

        story.append(Paragraph("Key Performance Indicators", styles["Heading2"]))
        story.append(Spacer(1, 6))

        kpi_data = []
        row = []
        palette = brand.palette or list(CHART_COLORS)
        for i, kpi in enumerate(kpis):
            value_style = ParagraphStyle(
                f"KPIValue{i}", parent=kpi_value,
                textColor=colors.HexColor(palette[i % len(palette)]),
            )
            cell = [Paragraph(kpi.name, kpi_label), Paragraph(kpi.formatted, value_style)]
            row.append(cell)
            if len(row) == 2 or i == len(kpis) - 1:
                if len(row) == 1:
                    row.append("")
                kpi_data.append(row)
                row = []

        kpi_table = Table(kpi_data, colWidths=[85 * mm, 85 * mm])
        kpi_table.setStyle(TableStyle([
            ("BACKGROUND", (0, 0), (-1, -1), colors.HexColor("#F9FAFB")),
            ("BOX", (0, 0), (-1, -1), 0.5, colors.HexColor("#E5E7EB")),
            ("INNERGRID", (0, 0), (-1, -1), 0.25, colors.HexColor("#E5E7EB")),
            ("TOPPADDING", (0, 0), (-1, -1), 10),
            ("BOTTOMPADDING", (0, 0), (-1, -1), 10),
            ("LEFTPADDING", (0, 0), (-1, -1), 12),
        ]))
        story.append(kpi_table)
        story.append(Spacer(1, 16))

        story.append(Paragraph("Executive Summary", styles["Heading2"]))
        story.append(Spacer(1, 6))
        for bullet in summary_bullets(summary_text, limit=6):
            story.append(Paragraph(f"• {bullet}", body_style))
            story.append(Spacer(1, 6))

        if charts:
            story.append(Spacer(1, 12))
            story.append(Paragraph("Visualisations", styles["Heading2"]))
            story.append(HRFlowable(width="100%", thickness=1, color=SECONDARY, spaceAfter=8))
            for chart_title, fig in charts[:6]:
                try:
                    styled = self._styled_figure(fig, brand)
                    img_bytes = self._figure_png(styled, width=900, height=450)
                    img_stream = io.BytesIO(img_bytes)
                    story.append(Paragraph(chart_title, styles["Heading3"]))
                    story.append(Spacer(1, 4))
                    story.append(Image(img_stream, width=170 * mm, height=85 * mm))
                    story.append(Spacer(1, 12))
                except Exception as e:
                    warning = f"Chart export failed for '{chart_title}': {e}"
                    logger.warning("PDF chart embed failed for '%s': %s", chart_title, e)
                    story.append(Paragraph(warning, body_style))
                    story.append(Spacer(1, 8))
        else:
            story.append(Spacer(1, 12))
            story.append(Paragraph("Visualisations", styles["Heading2"]))
            story.append(Spacer(1, 6))
            story.append(Paragraph("No visualisations were available for this report.", body_style))

        doc.build(story)
        logger.info("PDF exported: %s", path)
        return path

    # ── PowerPoint ─────────────────────────────────────────────────────────────
    def export_pptx(
        self,
        domain: DomainConfig,
        kpis: list[KPIResult],
        summary_text: str,
        charts: list[tuple[str, go.Figure]] | None = None,
        filename: str | None = None,
        brand: BrandTheme | None = None,
    ) -> Path:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

        brand = self._resolve_brand(brand)
        palette = brand.palette or list(CHART_COLORS)

        def rgb(hex_color: str) -> RGBColor:
            r, g, b = hex_to_rgb(hex_color)
            return RGBColor(r, g, b)

        BRAND = rgb(brand.primary)
        SECONDARY = rgb(brand.secondary)
        ACCENT = rgb(brand.accent)
        WHITE = RGBColor(0xFF, 0xFF, 0xFF)
        DARK = RGBColor(0x11, 0x18, 0x27)
        LIGHT = RGBColor(0xF9, 0xFA, 0xFB)
        MUTED = RGBColor(0x6B, 0x72, 0x80)

        filename = filename or f"deck_{domain.key}_{_timestamp()}.pptx"
        path = EXPORTS_DIR / filename

        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)
        blank = prs.slide_layouts[6]

        def add_rect(slide, l, t, w, h, fill_rgb):
            shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
            shape.fill.solid()
            shape.fill.fore_color.rgb = fill_rgb
            shape.line.fill.background()
            return shape

        def add_text(
            slide, text, l, t, w, h, size=18, bold=False,
            color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
        ):
            tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.vertical_anchor = anchor
            p = tf.paragraphs[0]
            p.alignment = align
            run = p.add_run()
            run.text = text
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = color
            return tb

        def add_header(slide, title: str, accent: RGBColor = BRAND):
            add_rect(slide, 0, 0, 13.33, 1.15, accent)
            add_rect(slide, 0, 1.15, 13.33, 0.08, SECONDARY)
            text_left = 1.6 if brand.logo_bytes else 0.45
            add_text(slide, title, text_left, 0.12, 11.5, 0.9, size=24, bold=True, color=WHITE)
            if brand.logo_bytes:
                logo_stream = io.BytesIO(brand.logo_bytes)
                slide.shapes.add_picture(logo_stream, Inches(0.35), Inches(0.18), height=Inches(0.75))

        # ── Slide 1: Cover ─────────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, BRAND)
        add_rect(slide, 8.8, 0, 4.53, 7.5, SECONDARY)
        add_rect(slide, 0, 6.55, 13.33, 0.95, ACCENT)

        if brand.logo_bytes:
            logo_stream = io.BytesIO(brand.logo_bytes)
            slide.shapes.add_picture(logo_stream, Inches(0.7), Inches(0.55), height=Inches(1.1))

        cover_title = brand.company_name or domain.label
        add_text(
            slide, cover_title, 0.7, 2.0, 7.5, 1.4,
            size=40, bold=True, color=WHITE,
        )
        add_text(slide, "Analytics Report", 0.7, 3.2, 7.5, 0.8, size=26, color=rgb("#E0E7FF"))
        add_text(
            slide,
            datetime.now().strftime("Generated %B %d, %Y"),
            0.7, 6.65, 8, 0.5, size=13, color=WHITE,
        )
        add_text(
            slide, "KPIs · Insights · Visualisations",
            9.1, 2.4, 3.8, 2.5, size=18, bold=True, color=WHITE, align=PP_ALIGN.LEFT,
        )

        # ── Slide 2: Agenda ────────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        add_header(slide, "What's Inside")
        agenda_items = [
            "Key performance indicators at a glance",
            "Executive summary and business insights",
            "Trend and category visualisations",
            f"Domain focus: {domain.label}",
        ]
        for idx, item in enumerate(agenda_items):
            y = 1.55 + idx * 1.15
            add_rect(slide, 0.5, y, 0.18, 0.55, rgb(palette[idx % len(palette)]))
            add_text(slide, item, 0.85, y - 0.05, 11.8, 0.7, size=20, color=DARK)

        # ── Slide 3: KPIs ──────────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Key Performance Indicators")

        cols = 4
        card_w, card_h = 2.85, 1.55
        x_start, y_start, gap = 0.35, 1.55, 0.28

        for i, kpi in enumerate(kpis[:8]):
            row, col = divmod(i, cols)
            x = x_start + col * (card_w + gap)
            y = y_start + row * (card_h + gap)
            card_color = rgb(palette[i % len(palette)])
            add_rect(slide, x, y, card_w, 0.12, card_color)
            card = add_rect(slide, x, y + 0.12, card_w, card_h - 0.12, LIGHT)
            card.line.color.rgb = rgb("#E5E7EB")
            add_text(slide, kpi.name, x + 0.12, y + 0.22, card_w - 0.24, 0.45, size=10, color=MUTED)
            add_text(
                slide, kpi.formatted, x + 0.12, y + 0.72, card_w - 0.24, 0.75,
                size=24, bold=True, color=card_color,
            )

        # ── Slide 4: Key Insights ─────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Key Insights", ACCENT)
        bullets = summary_bullets(summary_text, limit=5)
        for idx, bullet in enumerate(bullets):
            y = 1.55 + idx * 1.05
            chip = rgb(palette[idx % len(palette)])
            add_rect(slide, 0.45, y, 0.12, 0.75, chip)
            add_text(slide, bullet, 0.75, y, 11.9, 0.85, size=15, color=DARK)

        # ── Slide 5: Executive Summary ─────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        add_header(slide, "Executive Summary", SECONDARY)
        clean_summary = summary_text.replace("**", "").replace("##", "")
        add_text(slide, clean_summary[:1400], 0.45, 1.45, 12.4, 5.7, size=13, color=DARK)

        # ── Slide 6+: Charts ───────────────────────────────────────────────────
        if charts:
            for idx, (chart_title, fig) in enumerate(charts[:6]):
                slide = prs.slides.add_slide(blank)
                header_color = rgb(palette[idx % len(palette)])
                add_header(slide, chart_title, header_color)
                add_rect(slide, 0.35, 1.35, 12.6, 5.95, LIGHT)
                add_rect(slide, 0.35, 1.35, 0.12, 5.95, header_color)

                try:
                    styled = self._styled_figure(fig, brand)
                    img_bytes = self._figure_png(styled, width=1100, height=550)
                    img_stream = io.BytesIO(img_bytes)
                    slide.shapes.add_picture(
                        img_stream, Inches(0.55), Inches(1.5), Inches(12.2), Inches(5.65),
                    )
                except Exception as e:
                    logger.warning("Chart image export failed for '%s': %s", chart_title, e)
                    add_text(slide, f"[Chart unavailable: {chart_title}]", 0.55, 3.6, 12, 1, size=16, color=DARK)

        # ── Closing slide ────────────────────────────────────────────────────────
        slide = prs.slides.add_slide(blank)
        add_rect(slide, 0, 0, 13.33, 7.5, BRAND)
        add_rect(slide, 0, 6.2, 13.33, 1.3, ACCENT)
        if brand.logo_bytes:
            logo_stream = io.BytesIO(brand.logo_bytes)
            slide.shapes.add_picture(logo_stream, Inches(5.8), Inches(1.0), height=Inches(1.2))
        add_text(
            slide, "Thank You", 0.8, 2.8, 11.7, 1.0,
            size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            brand.company_name or domain.label,
            0.8, 3.9, 11.7, 0.7,
            size=20, color=rgb("#E0E7FF"), align=PP_ALIGN.CENTER,
        )
        add_text(
            slide,
            "Questions? Review the full dashboard or export additional formats.",
            0.8, 6.45, 11.7, 0.6,
            size=13, color=WHITE, align=PP_ALIGN.CENTER,
        )

        prs.save(str(path))
        logger.info("PPTX exported: %s", path)
        return path
