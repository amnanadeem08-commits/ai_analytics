"""
create_app_ppt.py — A short, focused overview deck for the AI Analytics
Assistant covering three pillars: App Intro · How It Works · Results.

This reuses the dark theme + reusable shape helpers from `create_ppt.py`
so the styling matches the full 12-slide deck. Run independently:

    python presentation/create_app_ppt.py

Output:
    presentation/app_overview.pptx   (6 slides)
"""

from __future__ import annotations

import sys
from pathlib import Path

# Reuse the theme + low-level helpers from the main generator.
sys.path.insert(0, str(Path(__file__).resolve().parent))
from create_ppt import (  # noqa: E402
    BG, PANEL, PANEL_2, BORDER, ACCENT, VIOLET, TEAL, TEXT, MUTED, FAINT,
    SUCCESS, WHITE, BODY_FONT, HEAD_FONT, MONO_FONT, SLIDE_W, SLIDE_H,
    slide, text, bullets, rounded, shape_text, chip, arrow, accent_bar,
    glow, notes, heading, _style_run,
)

from pptx import Presentation  # noqa: E402
from pptx.util import Inches, Pt  # noqa: E402
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # noqa: E402
from pptx.enum.shapes import MSO_SHAPE  # noqa: E402


APP_NAME = "AI Analytics Assistant"
APP_VER = "v3.0"
HERE = Path(__file__).resolve().parent
OUT_PATH = HERE / "app_overview.pptx"
TOTAL = 6


def footer(s, idx: int) -> None:
    """Divider + app tag + slide number (N / TOTAL)."""
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.02),
                            Inches(12.13), Pt(1))
    ln.shadow.inherit = False
    ln.fill.solid(); ln.fill.fore_color.rgb = BORDER
    ln.line.fill.background()
    text(s, f"{APP_NAME}  ·  {APP_VER}", Inches(0.6), Inches(7.06),
         Inches(8), Inches(0.3), size=9.5, color=FAINT)
    text(s, f"{idx:02d} / {TOTAL:02d}", Inches(11.2), Inches(7.06), Inches(1.53),
         Inches(0.3), size=9.5, color=FAINT, align=PP_ALIGN.RIGHT, font=MONO_FONT)


def section_tag(s, label: str) -> None:
    """A small rounded section badge top-right of a content slide."""
    chip(s, label, Inches(10.5), Inches(0.62), w=Inches(2.23), h=Inches(0.4),
         text_color=TEAL, border=BORDER, size=11)


# ─────────────────────────────────────────────────────────────────────────────
# Slides
# ─────────────────────────────────────────────────────────────────────────────

def s1_title(prs):
    s = slide(prs)
    glow(s, Inches(1.4), Inches(0.4), Inches(2.4), PANEL)
    glow(s, Inches(12.6), Inches(7.4), Inches(2.6), PANEL)
    accent_bar(s)

    tile = rounded(s, Inches(0.9), Inches(1.9), Inches(0.95), Inches(0.95),
                   fill=ACCENT, line=None, radius=0.22)
    shape_text(tile, [{"t": "📊", "size": 30, "color": WHITE, "align": PP_ALIGN.CENTER}],
               anchor=MSO_ANCHOR.MIDDLE)
    chip(s, "APP OVERVIEW · INTRO · WORKING · RESULTS", Inches(2.05), Inches(2.08),
         w=Inches(4.3), h=Inches(0.4), size=10.5)

    text(s, APP_NAME, Inches(0.9), Inches(3.0), Inches(11.5), Inches(1.3),
         size=50, color=TEXT, bold=True, font=HEAD_FONT, line_spacing=1.0)
    text(s, "Upload any spreadsheet → instant, domain-aware insights.",
         Inches(0.92), Inches(4.35), Inches(11.3), Inches(0.6),
         size=19, color=VIOLET, bold=True)
    text(s,
         "An AI analytics workspace that cleans your data, detects the business domain, "
         "computes KPIs, builds interactive Power BI-grade charts, writes an executive "
         "summary, answers questions in natural language, and exports branded reports.",
         Inches(0.92), Inches(4.95), Inches(11.2), Inches(1.2),
         size=14.5, color=MUTED, line_spacing=1.3)

    for i, t in enumerate(["Streamlit + Plotly", "10 Domains", "NL→SQL · RAG", "PDF · PPTX · Excel"]):
        chip(s, t, Inches(0.92) + Inches(2.62) * i, Inches(6.25),
             w=Inches(2.45), h=Inches(0.42), text_color=TEAL, border=BORDER, size=11)

    notes(s,
        "This is the AI Analytics Assistant. The one-line promise: drop in any CSV or Excel "
        "file and get analyst-grade output in seconds. In this short overview I'll cover "
        "three things — what the app is, how it works end to end, and what results you get back.")
    footer(s, 1)


def s2_intro(prs):
    s = slide(prs)
    heading(s, "App Intro", "What it is, in one screen")
    section_tag(s, "INTRO")

    band = rounded(s, Inches(0.6), Inches(1.95), Inches(12.13), Inches(1.5),
                   fill=PANEL_2, line=BORDER)
    shape_text(band, [
        {"t": "A senior data analyst, as an app.", "size": 16, "color": TEAL, "bold": True, "space_after": 6},
        {"t": "It replaces the slow, manual chain of clean → chart → analyze → write-up with a single "
              "upload. One domain-aware engine adapts the KPIs, charts, and AI focus to your data — "
              "and explains the findings instead of just plotting them.",
         "size": 13.5, "color": TEXT, "line_spacing": 1.3},
    ], anchor=MSO_ANCHOR.TOP, margin=0.28)

    cards = [
        ("🧠", "Domain-aware", "Detects 10 business domains with a confidence score and tailors insights."),
        ("📊", "Interactive BI", "Power BI-grade charts, KPI cross-filtering, one-screen Executive View."),
        ("💬", "Conversational", "NL chat, NL→SQL on DuckDB, and a lightweight RAG insight engine."),
        ("📤", "Deliver-ready", "Branded XLSX / PDF / PPTX exports — runs even with no AI key."),
    ]
    x0, y = Inches(0.6), Inches(3.7)
    cw, ch, gx = Inches(2.95), Inches(2.7), Inches(0.18)
    for i, (icon, title, body) in enumerate(cards):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        shape_text(card, [
            {"t": icon, "size": 26, "color": ACCENT, "align": PP_ALIGN.CENTER, "space_after": 8},
            {"t": title, "size": 15, "color": TEXT, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 6},
            {"t": body, "size": 11.5, "color": MUTED, "align": PP_ALIGN.CENTER, "line_spacing": 1.22},
        ], anchor=MSO_ANCHOR.TOP, margin=0.18)

    notes(s,
        "At its core the app is a senior data analyst packaged as software. It collapses the "
        "manual workflow into one upload. Four things define it: it's domain-aware with a "
        "confidence score, it's interactive BI with cross-filtering and an executive view, "
        "it's conversational through chat and NL-to-SQL, and it's deliver-ready with branded "
        "exports — and it still works without an AI key by falling back to statistics.")
    footer(s, 2)


def s3_working(prs):
    s = slide(prs)
    heading(s, "How It Works", "From upload to decision — one flow")
    section_tag(s, "WORKING")

    steps = [
        ("①", "Upload", "CSV / XLSX"),
        ("②", "Auto-Clean", "type-infer · dedupe"),
        ("③", "Detect Domain", "+ confidence"),
        ("④", "KPIs & Charts", "BI-grade visuals"),
        ("⑤", "AI Explain", "summary · story"),
        ("⑥", "Ask & Export", "NL→SQL · reports"),
    ]
    x0, y = Inches(0.6), Inches(2.2)
    cw, ch, gx = Inches(1.86), Inches(1.75), Inches(0.165)
    for i, (n, t, d) in enumerate(steps):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        shape_text(card, [
            {"t": n, "size": 20, "color": ACCENT, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 4},
            {"t": t, "size": 13, "color": TEXT, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 3},
            {"t": d, "size": 10, "color": MUTED, "align": PP_ALIGN.CENTER, "line_spacing": 1.1},
        ], anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            arrow(s, cx + cw - Inches(0.01), y + Inches(0.6), Inches(0.16), Inches(0.5),
                  color=BORDER)

    b1 = rounded(s, Inches(0.6), Inches(4.5), Inches(5.96), Inches(2.05),
                 fill=PANEL_2, line=BORDER)
    shape_text(b1, [
        {"t": "⚡  Fast path — on upload", "size": 14, "color": TEAL, "bold": True, "space_after": 7},
        {"t": "Cleaning, domain detection, profiling, KPIs, analytics and charts run "
              "instantly so the dashboard appears in seconds.",
         "size": 12.5, "color": TEXT, "line_spacing": 1.3},
    ], anchor=MSO_ANCHOR.TOP, margin=0.26)

    b2 = rounded(s, Inches(6.77), Inches(4.5), Inches(5.96), Inches(2.05),
                 fill=PANEL_2, line=BORDER)
    shape_text(b2, [
        {"t": "🧠  On demand — per tab (cached)", "size": 14, "color": VIOLET, "bold": True, "space_after": 7},
        {"t": "Expensive LLM work — executive summary, data story, insights, NL→SQL and "
              "exports — runs lazily only when you open that tab.",
         "size": 12.5, "color": TEXT, "line_spacing": 1.3},
    ], anchor=MSO_ANCHOR.TOP, margin=0.26)

    notes(s,
        "The flow is six steps: upload, auto-clean, detect the domain with a confidence "
        "score, compute KPIs and BI-grade charts, let the AI explain the findings, then ask "
        "questions or export. Under the hood it's two-tier: the fast deterministic path runs "
        "on upload so the dashboard is instant, while expensive AI work runs lazily per tab "
        "and is cached — keeping the app responsive and easy on API limits.")
    footer(s, 3)


def s4_results(prs):
    s = slide(prs)
    heading(s, "Results", "What you get back — in seconds")
    section_tag(s, "RESULTS")

    metrics = [("Seconds", "to first dashboard", ACCENT),
               ("10", "business domains", TEAL),
               ("3", "export formats", VIOLET),
               ("0", "API key needed to run", SUCCESS)]
    x0, y = Inches(0.6), Inches(2.0)
    cw, ch, gx = Inches(2.95), Inches(1.5), Inches(0.18)
    for i, (val, lab, c) in enumerate(metrics):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        shape_text(card, [
            {"t": val, "size": 30, "color": c, "bold": True, "font": MONO_FONT, "align": PP_ALIGN.CENTER, "space_after": 2},
            {"t": lab, "size": 11, "color": MUTED, "align": PP_ALIGN.CENTER, "line_spacing": 1.1},
        ], anchor=MSO_ANCHOR.MIDDLE)

    band = rounded(s, Inches(0.6), Inches(3.85), Inches(12.13), Inches(2.7),
                   fill=PANEL_2, line=BORDER)
    shape_text(band, [
        {"t": "Every upload produces", "size": 15, "color": TEAL, "bold": True, "space_after": 8},
        {"t": "•  Interactive dashboard — KPI cards, smart charts, data-quality score, KPI→chart cross-filtering & a one-screen Executive View.",
         "size": 13, "color": TEXT, "line_spacing": 1.3, "space_after": 4},
        {"t": "•  AI narrative — executive summary, chart-by-chart data story, and ranked business insights with recommendations.",
         "size": 13, "color": TEXT, "line_spacing": 1.3, "space_after": 4},
        {"t": "•  Answers on demand — ask in plain English; NL→SQL runs on DuckDB and a light RAG engine returns explained results.",
         "size": 13, "color": TEXT, "line_spacing": 1.3, "space_after": 4},
        {"t": "•  Branded deliverables — export a logo-and-palette PDF, PPTX or Excel report ready to share.",
         "size": 13, "color": TEXT, "line_spacing": 1.3},
    ], anchor=MSO_ANCHOR.TOP, margin=0.28)

    notes(s,
        "The result: a dashboard in seconds, fully domain-aware. You get an interactive "
        "dashboard with KPI cross-filtering and an executive view, an AI narrative — summary, "
        "data story and ranked insights — answers to plain-English questions via NL-to-SQL, "
        "and branded PDF/PPTX/Excel exports. It runs locally and deploys free on Streamlit Cloud.")
    footer(s, 4)


def s5_snapshot(prs):
    s = slide(prs)
    heading(s, "At a Glance", "Intro · Working · Results — summarized")
    section_tag(s, "RECAP")

    cols = [
        ("App Intro", ACCENT,
         ["Senior analyst, as an app",
          "Domain-aware across 10 verticals",
          "Explains findings, not just charts"]),
        ("How It Works", VIOLET,
         ["Upload → clean → detect domain",
          "KPIs + BI-grade charts instantly",
          "AI + NL→SQL run on demand"]),
        ("Results", TEAL,
         ["Dashboard in seconds",
          "Cross-filter + Executive View",
          "Branded PDF / PPTX / Excel"]),
    ]
    x0, y = Inches(0.6), Inches(2.05)
    cw, ch, gx = Inches(3.92), Inches(4.4), Inches(0.18)
    for i, (title, c, items) in enumerate(cols):
        cx = x0 + (cw + gx) * i
        rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        top = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, cw, Inches(0.5))
        top.shadow.inherit = False
        top.fill.solid(); top.fill.fore_color.rgb = c
        top.line.fill.background()
        try: top.adjustments[0] = 0.12
        except Exception: pass
        shape_text(top, [{"t": title, "size": 14, "color": WHITE, "bold": True,
                          "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        body = s.shapes.add_textbox(cx + Inches(0.22), y + Inches(0.75),
                                    cw - Inches(0.44), ch - Inches(0.9))
        tf = body.text_frame; tf.word_wrap = True
        for j, it in enumerate(items):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = 1.25; p.space_after = Pt(10)
            r = p.add_run(); r.text = f"•  {it}"
            _style_run(r, 12.5, TEXT, False, BODY_FONT)

    notes(s,
        "To summarize the three pillars: the intro — a domain-aware analyst app that explains, "
        "not just plots. The working — upload, clean, detect, chart instantly, with AI and "
        "NL-to-SQL on demand. The results — a dashboard in seconds with cross-filtering, an "
        "executive view, and branded exports.")
    footer(s, 5)


def s6_closing(prs):
    s = slide(prs)
    glow(s, Inches(11.9), Inches(0.6), Inches(2.4), PANEL)
    glow(s, Inches(1.2), Inches(7.2), Inches(2.4), PANEL)
    accent_bar(s)

    tile = rounded(s, Inches(0.9), Inches(2.2), Inches(0.95), Inches(0.95),
                   fill=ACCENT, line=None, radius=0.22)
    shape_text(tile, [{"t": "📊", "size": 30, "color": WHITE, "align": PP_ALIGN.CENTER}],
               anchor=MSO_ANCHOR.MIDDLE)

    text(s, "Thank You", Inches(0.9), Inches(3.35), Inches(11), Inches(1.1),
         size=48, color=TEXT, bold=True, font=HEAD_FONT)
    text(s, "From raw data to a decision — automatically, in one workspace.",
         Inches(0.92), Inches(4.55), Inches(11), Inches(0.6),
         size=18, color=VIOLET, bold=True)
    chip(s, "▶  streamlit run app.py  ·  localhost:8501", Inches(0.92), Inches(5.35),
         w=Inches(4.7), h=Inches(0.45), text_color=TEAL, border=BORDER, size=12)

    notes(s,
        "Thank you. That's the app in three pillars — intro, working, and results. It's live "
        "with 'streamlit run app.py' at localhost:8501 and deploys on Streamlit Cloud. Happy "
        "to give a live walkthrough on a real dataset.")
    footer(s, 6)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    s1_title(prs)
    s2_intro(prs)
    s3_working(prs)
    s4_results(prs)
    s5_snapshot(prs)
    s6_closing(prs)
    prs.save(str(OUT_PATH))
    return OUT_PATH


if __name__ == "__main__":
    out = build()
    print(f"Presentation generated: {out}")
    print(f"Slides: {TOTAL}  ·  Size: {out.stat().st_size / 1024:.0f} KB")
