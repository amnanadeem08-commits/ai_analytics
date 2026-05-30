"""
create_ppt.py — Auto-generates a professional presentation for the
AI Data Reporting & Analytics Assistant project.

Usage:
    python presentation/create_ppt.py

Output:
    presentation/project_presentation.pptx

Design notes
------------
- 16:9 widescreen, dark "command-center" theme matching the product UI.
- Reusable helpers (`slide`, `text`, `bullets`, `rounded`, `chip`, `arrow`,
  `footer`, `notes`) so the deck is easy to extend or restyle.
- Architecture (Slide 4) and Workflow (Slide 6) are drawn from native
  PowerPoint shapes, so they remain fully editable.
- Screenshots (Slide 8) auto-embed any images found in
  `presentation/screenshots/`; otherwise labeled placeholders are drawn.
- Content is taken from the actual codebase, README, and roadmap — not generic.

All content lives in builder functions; tweak text/colors in one place.
"""

from __future__ import annotations

import time
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ─────────────────────────────────────────────────────────────────────────────
# Theme
# ─────────────────────────────────────────────────────────────────────────────

BG       = RGBColor(0x0B, 0x10, 0x20)   # deep navy canvas
PANEL    = RGBColor(0x15, 0x1B, 0x2E)   # surface
PANEL_2  = RGBColor(0x1B, 0x23, 0x3A)   # lighter surface
BORDER   = RGBColor(0x2A, 0x33, 0x4E)   # hairline
ACCENT   = RGBColor(0x63, 0x66, 0xF1)   # indigo
VIOLET   = RGBColor(0xA8, 0x55, 0xF7)   # violet
TEAL     = RGBColor(0x38, 0xBD, 0xF8)   # sky
TEXT     = RGBColor(0xEC, 0xEE, 0xF6)   # primary text
MUTED    = RGBColor(0x9A, 0xA4, 0xBC)   # muted text
FAINT    = RGBColor(0x6B, 0x74, 0x90)   # faint
SUCCESS  = RGBColor(0x34, 0xD3, 0x99)
WARNING  = RGBColor(0xFB, 0xBF, 0x24)
ERROR    = RGBColor(0xF8, 0x71, 0x71)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)

BODY_FONT = "Segoe UI"
HEAD_FONT = "Segoe UI Semibold"
MONO_FONT = "Consolas"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PROJECT_NAME = "AI Data Reporting & Analytics Assistant"
PROJECT_VER  = "v3.0"

HERE = Path(__file__).resolve().parent
SHOTS_DIR = HERE / "screenshots"
OUT_PATH = HERE / "project_presentation.pptx"


# ─────────────────────────────────────────────────────────────────────────────
# Low-level helpers
# ─────────────────────────────────────────────────────────────────────────────

def slide(prs: Presentation):
    """Add a blank slide with the dark background applied."""
    s = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    fill = s.background.fill
    fill.solid()
    fill.fore_color.rgb = BG
    return s


def _style_run(run, size, color, bold, font, italic=False):
    f = run.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font


def text(s, content, x, y, w, h, *, size=18, color=TEXT, bold=False,
         font=BODY_FONT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         line_spacing=1.05, italic=False):
    """Single-paragraph textbox."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = content
    _style_run(run, size, color, bold, font, italic)
    return box


def bullets(s, items, x, y, w, h, *, size=15, color=TEXT, gap=8,
            marker="—", marker_color=ACCENT, line_spacing=1.08):
    """Multi-line bullet list. `items` may be str or (lead, rest) tuples."""
    box = s.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after = Pt(gap)
        m = p.add_run()
        m.text = f"{marker}  "
        _style_run(m, size, marker_color, True, BODY_FONT)
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run(); r1.text = lead
            _style_run(r1, size, TEXT, True, BODY_FONT)
            r2 = p.add_run(); r2.text = rest
            _style_run(r2, size, MUTED, False, BODY_FONT)
        else:
            r = p.add_run(); r.text = item
            _style_run(r, size, color, False, BODY_FONT)
    return box


def rounded(s, x, y, w, h, *, fill=PANEL, line=BORDER, line_w=Pt(1),
            radius=0.06, shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = line_w
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    return sp


def shape_text(sp, lines, *, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE,
               margin=0.12):
    """Write formatted lines into a shape. `lines` = list of dicts."""
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(margin)
    tf.margin_right = Inches(margin)
    tf.margin_top = Inches(0.06)
    tf.margin_bottom = Inches(0.06)
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = ln.get("align", align)
        p.line_spacing = ln.get("line_spacing", 1.0)
        if "space_after" in ln:
            p.space_after = Pt(ln["space_after"])
        run = p.add_run()
        run.text = ln["t"]
        _style_run(run, ln.get("size", 13), ln.get("color", TEXT),
                   ln.get("bold", False), ln.get("font", BODY_FONT))
    return sp


def chip(s, label, x, y, *, w=Inches(1.7), h=Inches(0.34), fill=None,
         text_color=ACCENT, border=ACCENT, size=10.5):
    sp = rounded(s, x, y, w, h, fill=fill, line=border, line_w=Pt(1), radius=0.5)
    shape_text(sp, [{"t": label, "size": size, "color": text_color,
                     "bold": True, "align": PP_ALIGN.CENTER}],
               anchor=MSO_ANCHOR.MIDDLE, margin=0.05)
    return sp


def arrow(s, x, y, w, h, *, color=ACCENT, shape=MSO_SHAPE.CHEVRON):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp


def kicker(s, label, x, y, *, color=ACCENT):
    """Small uppercase section label."""
    return text(s, label.upper(), x, y, Inches(8), Inches(0.4),
                size=12.5, color=color, bold=True, font=HEAD_FONT)


def accent_bar(s):
    """Thin three-tone bar across the very top of the slide."""
    third = SLIDE_W // 3
    for i, c in enumerate((ACCENT, VIOLET, TEAL)):
        b = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, third * i, 0, third, Inches(0.08))
        b.shadow.inherit = False
        b.fill.solid(); b.fill.fore_color.rgb = c
        b.line.fill.background()


def footer(s, idx):
    """Divider + project tag + slide number."""
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(7.02),
                            Inches(12.13), Pt(1))
    ln.shadow.inherit = False
    ln.fill.solid(); ln.fill.fore_color.rgb = BORDER
    ln.line.fill.background()
    text(s, f"{PROJECT_NAME}  ·  {PROJECT_VER}", Inches(0.6), Inches(7.06),
         Inches(8), Inches(0.3), size=9.5, color=FAINT)
    text(s, f"{idx:02d} / 12", Inches(11.2), Inches(7.06), Inches(1.53),
         Inches(0.3), size=9.5, color=FAINT, align=PP_ALIGN.RIGHT, font=MONO_FONT)


def heading(s, title, sub=None):
    """Standard content-slide heading with accent bar + kicker."""
    accent_bar(s)
    text(s, title, Inches(0.6), Inches(0.5), Inches(12.1), Inches(0.9),
         size=33, color=TEXT, bold=True, font=HEAD_FONT)
    if sub:
        text(s, sub, Inches(0.62), Inches(1.32), Inches(12.1), Inches(0.5),
             size=14.5, color=MUTED)


def notes(s, txt):
    s.notes_slide.notes_text_frame.text = txt.strip()


def glow(s, cx, cy, r, color):
    """Subtle radial-ish glow using a slightly-lighter oval (no transparency)."""
    sp = s.shapes.add_shape(MSO_SHAPE.OVAL, cx - r, cy - r, r * 2, r * 2)
    sp.shadow.inherit = False
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    sp.line.fill.background()
    return sp


# ─────────────────────────────────────────────────────────────────────────────
# Slide builders
# ─────────────────────────────────────────────────────────────────────────────

def slide_01_title(prs):
    s = slide(prs)
    # ambient glows
    glow(s, Inches(1.4), Inches(0.4), Inches(2.4), PANEL)
    glow(s, Inches(12.6), Inches(7.4), Inches(2.6), PANEL)
    accent_bar(s)

    # logo tile
    tile = rounded(s, Inches(0.9), Inches(2.0), Inches(0.95), Inches(0.95),
                   fill=ACCENT, line=None, radius=0.22)
    shape_text(tile, [{"t": "📊", "size": 30, "color": WHITE,
                       "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)

    chip(s, "AI · DATA ANALYTICS · BI", Inches(2.05), Inches(2.18),
         w=Inches(2.9), h=Inches(0.38))

    text(s, PROJECT_NAME, Inches(0.9), Inches(3.05), Inches(11.5), Inches(1.6),
         size=46, color=TEXT, bold=True, font=HEAD_FONT, line_spacing=1.0)
    text(s, "From raw spreadsheet to decision-ready insight — automatically.",
         Inches(0.92), Inches(4.55), Inches(11.3), Inches(0.6),
         size=19, color=VIOLET, bold=True)
    text(s,
         "An AI-powered analytics platform that cleans data, detects your business "
         "domain, computes KPIs, builds visuals, writes an executive summary, answers "
         "questions in natural language, and exports branded reports.",
         Inches(0.92), Inches(5.15), Inches(11.0), Inches(1.1),
         size=14.5, color=MUTED, line_spacing=1.25)

    text(s, "Your Name  ·  Data / AI Engineer", Inches(0.92), Inches(6.45),
         Inches(7), Inches(0.4), size=12.5, color=FAINT)
    notes(s,
        "Welcome. This is the AI Data Reporting & Analytics Assistant — a Streamlit "
        "app that acts like a senior data analyst. The core promise: upload any CSV or "
        "Excel file and get analyst-grade output in seconds, with zero manual setup. "
        "I'll cover the problem, the solution, architecture, tech stack, workflow, "
        "features, results, scalability, and business value.")
    footer(s, 1)


def slide_02_problem(prs):
    s = slide(prs)
    heading(s, "The Problem", "Why this project exists")
    kicker(s, "Pain points", Inches(0.6), Inches(1.95))

    cards = [
        ("⏳", "Slow & manual", "Turning a raw spreadsheet into clean KPIs, charts and a written summary takes a skilled analyst hours — per dataset."),
        ("🔒", "Analyst bottleneck", "Non-technical teams can't self-serve; every question waits in an analyst's queue."),
        ("🧩", "Fragmented tooling", "Cleaning, BI dashboards, stats, forecasting and report design live in 4–5 disconnected tools."),
        ("📉", "Insight, not just charts", "Most tools draw charts but don't explain what's happening, why, or what to do next."),
    ]
    x0, y0 = Inches(0.6), Inches(2.45)
    cw, ch, gx, gy = Inches(5.95), Inches(1.95), Inches(0.25), Inches(0.3)
    for i, (icon, title, body) in enumerate(cards):
        cx = x0 + (cw + gx) * (i % 2)
        cy = y0 + (ch + gy) * (i // 2)
        card = rounded(s, cx, cy, cw, ch, fill=PANEL, line=BORDER)
        shape_text(card, [
            {"t": f"{icon}  {title}", "size": 16, "color": TEXT, "bold": True,
             "space_after": 6},
            {"t": body, "size": 12.5, "color": MUTED, "line_spacing": 1.2},
        ], anchor=MSO_ANCHOR.TOP)
    notes(s,
        "The status quo is slow and manual. A competent analyst spends hours per file "
        "on cleaning, KPI definition, charting, and writing a narrative. Business teams "
        "can't self-serve, so they wait. The tooling is fragmented across spreadsheets, "
        "BI suites, notebooks and slide software. And critically — most tools stop at "
        "charts; they don't explain the 'so what'. That gap is what this project closes.")
    footer(s, 2)


def slide_03_solution(prs):
    s = slide(prs)
    heading(s, "The Solution",
            "One upload → a complete, domain-aware analytics workspace")
    kicker(s, "How it works at a glance", Inches(0.6), Inches(1.95))

    steps = [
        ("①", "Upload", "Drop a CSV / XLSX"),
        ("②", "Understand", "Auto-clean + detect domain"),
        ("③", "Analyze", "KPIs, charts, stats, anomalies"),
        ("④", "Explain", "AI summary, story, insights"),
        ("⑤", "Deliver", "Ask in NL · export reports"),
    ]
    x0, y = Inches(0.6), Inches(2.5)
    cw, ch, gx = Inches(2.28), Inches(1.7), Inches(0.18)
    for i, (n, t, d) in enumerate(steps):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        shape_text(card, [
            {"t": n, "size": 22, "color": ACCENT, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 4},
            {"t": t, "size": 15, "color": TEXT, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 3},
            {"t": d, "size": 11, "color": MUTED, "align": PP_ALIGN.CENTER, "line_spacing": 1.1},
        ], anchor=MSO_ANCHOR.MIDDLE)
        if i < len(steps) - 1:
            arrow(s, cx + cw + Inches(0.005), y + Inches(0.62),
                  Inches(0.17), Inches(0.45), color=BORDER)

    band = rounded(s, Inches(0.6), Inches(4.7), Inches(12.13), Inches(2.05),
                   fill=PANEL_2, line=BORDER)
    shape_text(band, [
        {"t": "What makes it different", "size": 15, "color": TEAL, "bold": True, "space_after": 6},
        {"t": "Domain-aware:  a Smart Intelligence layer detects your domain with a confidence score and filters out irrelevant framing.",
         "size": 12.5, "color": TEXT, "line_spacing": 1.22, "space_after": 3},
        {"t": "Explains, not just plots:  executive summary + data story answer what / why / what next.",
         "size": 12.5, "color": TEXT, "line_spacing": 1.22, "space_after": 3},
        {"t": "Interactive BI:  Power BI-grade charts, click-a-KPI cross-filtering, and a one-screen Executive View.",
         "size": 12.5, "color": TEXT, "line_spacing": 1.22, "space_after": 3},
        {"t": "Conversational:  natural-language chat, NL→SQL on DuckDB, and a lightweight RAG insight engine.",
         "size": 12.5, "color": TEXT, "line_spacing": 1.22, "space_after": 3},
        {"t": "Graceful by design:  works even with no AI key — falls back to statistical narratives.",
         "size": 12.5, "color": TEXT, "line_spacing": 1.22},
    ], anchor=MSO_ANCHOR.TOP)
    notes(s,
        "The solution collapses the whole analyst workflow into a single upload. Five "
        "stages: understand, analyze, explain, and deliver. The differentiators are: it's "
        "domain-aware across 10 verticals; it explains findings rather than just plotting "
        "them; it's conversational via chat, NL-to-SQL on DuckDB, and a light RAG engine; "
        "and it degrades gracefully to statistical summaries when no AI key is present.")
    footer(s, 3)


def slide_04_architecture(prs):
    s = slide(prs)
    heading(s, "Project Architecture",
            "Layered, decoupled design — UI never touches business logic; data flows top-down")

    layers = [
        ("Presentation · Streamlit UI", "app.py · 8 tabs · components · dark glass theme", ACCENT),
        ("Orchestration", "pipeline.py — run_pipeline · lazy ensure_* · filtered cache · cached services", VIOLET),
        ("Analytics Services", "cleaning · domain · profiling · KPI · analytics · charts · quality · forecasting", TEAL),
        ("AI Engines", "AI Summary · Storytelling · Insights + Recommendations · NL→SQL (DuckDB) · Light RAG (FAISS)", ACCENT),
        ("LLM Gateway", "utils/llm_client — Anthropic · OpenAI · OpenRouter  (fallback · retries · graceful degrade)", VIOLET),
        ("Storage & Output", "in-memory DuckDB · FAISS / NumPy vector store · session_state · exports → XLSX / PDF / PPTX", TEAL),
    ]
    x = Inches(0.85)
    w = Inches(11.6)
    y = Inches(1.95)
    bh = Inches(0.72)
    gap = Inches(0.115)
    for i, (title, comps, c) in enumerate(layers):
        cy = y + (bh + gap) * i
        band = rounded(s, x, cy, w, bh, fill=PANEL, line=BORDER)
        # accent tab on the left
        tab = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, cy, Inches(0.12), bh)
        tab.shadow.inherit = False
        tab.fill.solid(); tab.fill.fore_color.rgb = c
        tab.line.fill.background()
        try: tab.adjustments[0] = 0.5
        except Exception: pass
        shape_text(band, [
            {"t": title, "size": 14.5, "color": TEXT, "bold": True, "space_after": 2},
            {"t": comps, "size": 11.5, "color": MUTED, "line_spacing": 1.05},
        ], anchor=MSO_ANCHOR.MIDDLE, margin=0.28)
        if i < len(layers) - 1:
            dn = arrow(s, x + w/2 - Inches(0.11), cy + bh - Inches(0.02),
                       Inches(0.22), gap + Inches(0.04),
                       color=BORDER, shape=MSO_SHAPE.DOWN_ARROW)
    notes(s,
        "Architecture is strictly layered. The Streamlit UI dispatches to tabs but holds "
        "no business logic. pipeline.py orchestrates everything: it runs the fast, "
        "deterministic work eagerly and defers expensive LLM work behind lazy ensure_* "
        "getters, all cached. Below that sit pure analytics services, then the AI engines "
        "including NL-to-SQL on DuckDB and a FAISS-backed light RAG. A single LLM gateway "
        "abstracts three providers with fallback and retries. The bottom layer is storage "
        "and output. Data flows top-to-bottom; each layer is independently testable.")
    footer(s, 4)


def slide_05_stack(prs):
    s = slide(prs)
    heading(s, "Technical Stack", "Production-grade Python data & AI tooling")

    groups = [
        ("App & UI", ACCENT, ["Streamlit 1.35", "Python 3.11", "Custom dark CSS theme"]),
        ("Data & Stats", TEAL, ["Pandas 2.2", "NumPy 1.26", "SciPy 1.13", "scikit-learn 1.5"]),
        ("Visualization", VIOLET, ["Plotly 5.22", "Kaleido (static export)"]),
        ("AI / LLM", ACCENT, ["Anthropic Claude", "OpenAI GPT-4o", "OpenRouter (fallback)"]),
        ("SQL & Retrieval", TEAL, ["DuckDB 1.0 (NL→SQL)", "FAISS-cpu 1.8 (vectors)", "HashingVectorizer"]),
        ("Reporting & I/O", VIOLET, ["ReportLab (PDF)", "python-pptx (PPTX)", "openpyxl / XlsxWriter", "Pillow · python-dotenv"]),
    ]
    x0, y0 = Inches(0.6), Inches(2.05)
    cw, ch, gx, gy = Inches(3.92), Inches(2.15), Inches(0.18), Inches(0.25)
    for i, (title, c, items) in enumerate(groups):
        cx = x0 + (cw + gx) * (i % 3)
        cy = y0 + (ch + gy) * (i // 3)
        card = rounded(s, cx, cy, cw, ch, fill=PANEL, line=BORDER)
        lines = [{"t": title, "size": 14.5, "color": c, "bold": True, "space_after": 7}]
        for it in items:
            lines.append({"t": f"•  {it}", "size": 12, "color": TEXT,
                          "line_spacing": 1.15, "space_after": 3})
        shape_text(card, lines, anchor=MSO_ANCHOR.TOP)
    notes(s,
        "The stack is entirely Python 3.11. Streamlit drives the UI. Pandas, NumPy, SciPy "
        "and scikit-learn power data and statistics. Plotly plus Kaleido handle interactive "
        "and static charts. The AI layer supports Anthropic, OpenAI and OpenRouter with "
        "automatic fallback. DuckDB executes generated SQL in-memory; FAISS with a hashing "
        "vectorizer powers the light RAG. Reporting uses ReportLab, python-pptx and "
        "openpyxl. Everything is pinned in requirements.txt for reproducible deploys.")
    footer(s, 5)


def slide_06_workflow(prs):
    s = slide(prs)
    heading(s, "Workflow Process",
            "End-to-end flow — fast path runs on upload, AI work runs on demand")

    # Row 1: eager pipeline
    kicker(s, "Automatic on upload (fast path)", Inches(0.6), Inches(1.85), color=TEAL)
    eager = ["Upload\nCSV / XLSX", "Auto-Clean\n& type infer", "Detect\nDomain",
             "Profile · KPIs\n· Analytics", "Smart\nCharts"]
    x0, y = Inches(0.6), Inches(2.25)
    cw, ch, gx = Inches(2.18), Inches(1.25), Inches(0.22)
    for i, label in enumerate(eager):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        parts = label.split("\n")
        shape_text(card, [
            {"t": parts[0], "size": 13.5, "color": TEXT, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 2},
            {"t": parts[1] if len(parts) > 1 else "", "size": 10.5, "color": MUTED, "align": PP_ALIGN.CENTER},
        ], anchor=MSO_ANCHOR.MIDDLE)
        if i < len(eager) - 1:
            arrow(s, cx + cw - Inches(0.02), y + Inches(0.4), Inches(0.2), Inches(0.45),
                  color=TEAL)

    # connector down
    arrow(s, x0 + Inches(5.4), y + ch + Inches(0.02), Inches(0.3), Inches(0.4),
          color=BORDER, shape=MSO_SHAPE.DOWN_ARROW)

    # Row 2: on-demand AI
    kicker(s, "On demand, per tab (lazy + cached)", Inches(0.6), Inches(4.2), color=VIOLET)
    lazy = ["AI Executive\nSummary", "Data\nStory", "Insights &\nRecommendations",
            "NL → SQL\n(DuckDB)", "RAG\nQ&A", "Branded\nExport"]
    y2 = Inches(4.6)
    cw2, gx2 = Inches(1.85), Inches(0.17)
    for i, label in enumerate(lazy):
        cx = x0 + (cw2 + gx2) * i
        card = rounded(s, cx, y2, cw2, ch, fill=PANEL_2, line=BORDER)
        parts = label.split("\n")
        shape_text(card, [
            {"t": parts[0], "size": 12.5, "color": TEXT, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 2},
            {"t": parts[1] if len(parts) > 1 else "", "size": 10, "color": MUTED, "align": PP_ALIGN.CENTER},
        ], anchor=MSO_ANCHOR.MIDDLE)
        if i < len(lazy) - 1:
            arrow(s, cx + cw2 - Inches(0.02), y2 + Inches(0.4), Inches(0.18), Inches(0.45),
                  color=VIOLET)
    notes(s,
        "Two-tier workflow. On upload, only fast deterministic work runs: cleaning, domain "
        "detection, profiling, KPIs, analytics and charts — so the dashboard appears almost "
        "instantly. Expensive LLM work — the executive summary, data story, insights, "
        "NL-to-SQL, RAG answers and report export — runs lazily, only when the user opens "
        "that tab, and is cached in session state. This keeps the app responsive and avoids "
        "burning API rate limits. Filters recompute only the affected outputs via a signature-keyed cache.")
    footer(s, 6)


def slide_07_features(prs):
    s = slide(prs)
    heading(s, "Features & Functionality", "A cohesive, domain-aware BI workspace")

    col1 = [
        ("Auto Data Engine — ", "clean, dedupe, type-infer, time-series detect"),
        ("Smart Intelligence — ", "auto-detects 10 domains with a confidence score"),
        ("Insight Filtering — ", "drops irrelevant framing (no 'profit' on healthcare)"),
        ("KPI Engine — ", "domain-aware metrics, currency/% formatting & deltas"),
        ("BI-Grade Visuals — ", "Tableau / Power BI-quality interactive Plotly charts"),
    ]
    col2 = [
        ("Interactive Dashboard — ", "KPI→chart cross-filter + one-screen Executive View"),
        ("AI Summary & Story — ", "executive narrative + chart-by-chart storytelling"),
        ("Business Insights — ", "ranked findings, severity & recommendations"),
        ("Data Copilot — ", "NL chat · NL→SQL (DuckDB) · light RAG insights"),
        ("Branded Export — ", "logo + palette → XLSX / PDF / PPTX reports"),
    ]
    cw = Inches(6.0)
    c1 = rounded(s, Inches(0.6), Inches(2.0), cw, Inches(4.6), fill=PANEL, line=BORDER)
    c2 = rounded(s, Inches(6.73), Inches(2.0), cw, Inches(4.6), fill=PANEL, line=BORDER)
    bullets(s, col1, Inches(0.85), Inches(2.3), Inches(5.5), Inches(4.0), size=13.5, gap=14)
    bullets(s, col2, Inches(6.98), Inches(2.3), Inches(5.5), Inches(4.0), size=13.5, gap=14, marker_color=VIOLET)
    notes(s,
        "The data engine cleans and types data automatically. A Smart Intelligence layer "
        "detects the vertical with a confidence score and filters out domain-irrelevant "
        "framing. KPIs are domain-aware, and the visuals are Tableau/Power BI-grade Plotly "
        "charts. The dashboard is interactive: click a KPI to cross-filter every chart, or "
        "flip the Executive View toggle for a single-screen summary. On the AI side: "
        "executive summary and data story, ranked business insights with recommendations, "
        "and the Data Copilot — natural-language chat, NL-to-SQL on DuckDB, and a light RAG "
        "engine — plus fully branded multi-format export.")
    footer(s, 7)


def slide_08_screenshots(prs):
    s = slide(prs)
    heading(s, "Product Screenshots", "Dark command-center UI — drop images into presentation/screenshots/")

    labels = ["Dashboard — KPIs & Charts", "Data Copilot — NL→SQL",
              "Business Insights", "Branded Export"]
    imgs = []
    if SHOTS_DIR.exists():
        for ext in ("*.png", "*.jpg", "*.jpeg"):
            imgs.extend(sorted(SHOTS_DIR.glob(ext)))

    x0, y0 = Inches(0.6), Inches(2.0)
    cw, ch, gx, gy = Inches(5.95), Inches(2.2), Inches(0.23), Inches(0.28)
    for i in range(4):
        cx = x0 + (cw + gx) * (i % 2)
        cy = y0 + (ch + gy) * (i // 2)
        if i < len(imgs):
            frame = rounded(s, cx, cy, cw, ch, fill=PANEL, line=ACCENT)
            try:
                s.shapes.add_picture(str(imgs[i]), cx + Inches(0.08), cy + Inches(0.08),
                                     width=cw - Inches(0.16), height=ch - Inches(0.16))
            except Exception:
                pass
        else:
            ph = rounded(s, cx, cy, cw, ch, fill=PANEL, line=BORDER, line_w=Pt(1.25))
            shape_text(ph, [
                {"t": "🖼", "size": 26, "color": FAINT, "align": PP_ALIGN.CENTER, "space_after": 4},
                {"t": f"Screenshot placeholder", "size": 13, "color": MUTED, "bold": True, "align": PP_ALIGN.CENTER, "space_after": 2},
                {"t": labels[i], "size": 11.5, "color": ACCENT, "align": PP_ALIGN.CENTER},
            ], anchor=MSO_ANCHOR.MIDDLE)
    notes(s,
        "These are the key screens to capture: the main dashboard with KPI cards and "
        "charts, the Data Copilot showing a natural-language-to-SQL result, the Business "
        "Insights panel, and a generated branded export. To swap the placeholders for real "
        "images, drop PNG/JPG files into presentation/screenshots/ and re-run create_ppt.py — "
        "the first four images auto-fill these frames.")
    footer(s, 8)


def slide_09_results(prs):
    s = slide(prs)
    heading(s, "Results & Demo", "What the user gets back — in seconds")

    # metric tiles
    metrics = [("Seconds", "to first dashboard", ACCENT),
               ("10", "business domains supported", TEAL),
               ("3", "export formats (XLSX/PDF/PPTX)", VIOLET),
               ("0", "API key required to run", SUCCESS)]
    x0, y = Inches(0.6), Inches(2.0)
    cw, ch, gx = Inches(2.95), Inches(1.5), Inches(0.18)
    for i, (val, lab, c) in enumerate(metrics):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        shape_text(card, [
            {"t": val, "size": 30, "color": c, "bold": True, "font": MONO_FONT, "align": PP_ALIGN.CENTER, "space_after": 2},
            {"t": lab, "size": 11, "color": MUTED, "align": PP_ALIGN.CENTER, "line_spacing": 1.1},
        ], anchor=MSO_ANCHOR.MIDDLE)

    band = rounded(s, Inches(0.6), Inches(3.85), Inches(12.13), Inches(2.7),
                   fill=PANEL_2, line=BORDER)
    shape_text(band, [
        {"t": "Demo walkthrough", "size": 15, "color": TEAL, "bold": True, "space_after": 8},
        {"t": "1.  Upload a sales / healthcare / finance CSV  →  data is cleaned & domain auto-detected.",
         "size": 13, "color": TEXT, "line_spacing": 1.3, "space_after": 4},
        {"t": "2.  Dashboard renders KPI cards, BI-grade charts & a quality score — click a KPI to cross-filter, or flip to Executive View.",
         "size": 13, "color": TEXT, "line_spacing": 1.3, "space_after": 4},
        {"t": "3.  Open Data Story / Insights  →  AI explains what's happening, why, and what to do.",
         "size": 13, "color": TEXT, "line_spacing": 1.3, "space_after": 4},
        {"t": "4.  Ask in the Copilot: \"top 5 regions by revenue\"  →  NL→SQL runs on DuckDB & returns a table.",
         "size": 13, "color": TEXT, "line_spacing": 1.3, "space_after": 4},
        {"t": "5.  Export  →  a branded PDF / PPTX / Excel report with your logo and palette.",
         "size": 13, "color": TEXT, "line_spacing": 1.3},
    ], anchor=MSO_ANCHOR.TOP)
    notes(s,
        "In a live demo: upload a dataset and the dashboard appears in seconds, fully "
        "domain-aware. The data story and insights tabs generate AI narratives on demand. "
        "The Copilot turns a plain-English question into SQL, runs it on DuckDB, and returns "
        "a result table with an explanation. Finally, export produces a branded PDF, PPTX or "
        "Excel report. It runs locally at localhost:8501 and deploys free on Streamlit Cloud.")
    footer(s, 9)


def slide_10_scalability(prs):
    s = slide(prs)
    heading(s, "Scalability & Future Improvements", "A roadmap already scoped in the codebase (ROADMAP_V4)")

    cols = [
        ("Predictive Intelligence", ACCENT,
         ["ARIMA / Prophet forecasting + confidence bands",
          "Churn prediction & customer segmentation",
          "What-if scenario & sensitivity modeling"]),
        ("Exploration & Depth", TEAL,
         ["Interactive chart drill-down to rows",
          "Cohort analysis & retention curves",
          "Pivot tables & custom metric builder"]),
        ("Scale & Collaboration", VIOLET,
         ["Auth, multi-user workspaces & sharing",
          "Scheduled email / HTML / Power BI export",
          "Vector-DB & warehouse backends, async LLM"]),
    ]
    x0, y = Inches(0.6), Inches(2.05)
    cw, ch, gx = Inches(3.92), Inches(4.4), Inches(0.18)
    for i, (title, c, items) in enumerate(cols):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        top = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, cx, y, cw, Inches(0.5))
        top.shadow.inherit = False
        top.fill.solid(); top.fill.fore_color.rgb = c
        top.line.fill.background()
        try: top.adjustments[0] = 0.12
        except Exception: pass
        shape_text(top, [{"t": title, "size": 13.5, "color": WHITE, "bold": True,
                          "align": PP_ALIGN.CENTER}], anchor=MSO_ANCHOR.MIDDLE)
        lines = []
        for it in items:
            lines.append({"t": f"•  {it}", "size": 12.5, "color": TEXT,
                          "line_spacing": 1.25, "space_after": 9})
        body = s.shapes.add_textbox(cx + Inches(0.2), y + Inches(0.75),
                                    cw - Inches(0.4), ch - Inches(0.9))
        tf = body.text_frame; tf.word_wrap = True
        for j, ln in enumerate(lines):
            p = tf.paragraphs[0] if j == 0 else tf.add_paragraph()
            p.line_spacing = ln["line_spacing"]; p.space_after = Pt(ln["space_after"])
            r = p.add_run(); r.text = ln["t"]
            _style_run(r, ln["size"], ln["color"], False, BODY_FONT)
    notes(s,
        "Scalability is already planned in ROADMAP_V4. Phase 2 adds predictive intelligence "
        "— ARIMA/Prophet forecasting, churn and segmentation, and what-if scenarios. Phase 3 "
        "adds exploration depth: drill-down, cohorts and pivot tables. Phases 4–5 cover "
        "collaboration and scale: auth and multi-user workspaces, scheduled and Power BI "
        "exports, and swapping in a managed vector DB and warehouse backends with async LLM "
        "calls. The decoupled service architecture means each of these slots in cleanly.")
    footer(s, 10)


def slide_11_value(prs):
    s = slide(prs)
    heading(s, "Business, Client & Portfolio Value", "Why it matters to each audience")

    cards = [
        ("🏢", "Business value", ACCENT,
         ["Hours of analyst work → seconds",
          "Self-serve analytics for non-technical teams",
          "Consistent, domain-aware reporting"]),
        ("🤝", "Client value", TEAL,
         ["White-label branded deliverables",
          "Deployable & demo-ready in minutes",
          "Runs with or without an AI budget"]),
        ("💼", "Portfolio value", VIOLET,
         ["Full-stack AI + data engineering",
          "Clean layered architecture & caching",
          "LLM orchestration, RAG, NL→SQL, BI"]),
    ]
    x0, y = Inches(0.6), Inches(2.05)
    cw, ch, gx = Inches(3.92), Inches(4.4), Inches(0.18)
    for i, (icon, title, c, items) in enumerate(cards):
        cx = x0 + (cw + gx) * i
        card = rounded(s, cx, y, cw, ch, fill=PANEL, line=BORDER)
        lines = [{"t": f"{icon}", "size": 26, "color": c, "align": PP_ALIGN.LEFT, "space_after": 6},
                 {"t": title, "size": 16, "color": TEXT, "bold": True, "space_after": 10}]
        for it in items:
            lines.append({"t": f"•  {it}", "size": 12.5, "color": MUTED,
                          "line_spacing": 1.25, "space_after": 8})
        shape_text(card, lines, anchor=MSO_ANCHOR.TOP, margin=0.22)
    notes(s,
        "Three audiences, three angles. For the business: it converts hours of analyst time "
        "into seconds and lets non-technical teams self-serve. For a client: it produces "
        "white-label branded reports, deploys in minutes, and runs with or without an AI "
        "budget thanks to graceful fallback. For a portfolio: it demonstrates end-to-end "
        "AI and data engineering — clean layered architecture, performance caching, LLM "
        "orchestration with fallback, RAG, NL-to-SQL, and BI reporting in one product.")
    footer(s, 11)


def slide_12_closing(prs):
    s = slide(prs)
    glow(s, Inches(11.9), Inches(0.6), Inches(2.4), PANEL)
    glow(s, Inches(1.2), Inches(7.2), Inches(2.4), PANEL)
    accent_bar(s)

    tile = rounded(s, Inches(0.9), Inches(2.2), Inches(0.95), Inches(0.95),
                   fill=ACCENT, line=None, radius=0.22)
    shape_text(tile, [{"t": "📊", "size": 30, "color": WHITE, "align": PP_ALIGN.CENTER}],
               anchor=MSO_ANCHOR.MIDDLE)

    text(s, "Thank You", Inches(0.9), Inches(3.35), Inches(11), Inches(1.1),
         size=48, color=TEXT, bold=True, font=HEAD_FONT)
    text(s, "From raw data to a decision — automatically, in one workspace.",
         Inches(0.92), Inches(4.55), Inches(11), Inches(0.6),
         size=18, color=VIOLET, bold=True)

    chip(s, "▶  streamlit run app.py  ·  localhost:8501", Inches(0.92), Inches(5.35),
         w=Inches(4.7), h=Inches(0.45), text_color=TEAL, border=BORDER, size=12)
    text(s, "Your Name  ·  email@example.com  ·  github.com/your-handle",
         Inches(0.92), Inches(6.1), Inches(11), Inches(0.4),
         size=13, color=MUTED)
    notes(s,
        "Thank you. To recap: this turns raw data into a decision automatically, in a single "
        "workspace — clean, analyze, explain, converse, and export. It's live with "
        "'streamlit run app.py' at localhost:8501 and deploys on Streamlit Cloud. "
        "Happy to walk through the code, the architecture, or a live dataset. "
        "Replace the contact line with your details before sending.")
    footer(s, 12)


# ─────────────────────────────────────────────────────────────────────────────
# Build
# ─────────────────────────────────────────────────────────────────────────────

def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_problem(prs)
    slide_03_solution(prs)
    slide_04_architecture(prs)
    slide_05_stack(prs)
    slide_06_workflow(prs)
    slide_07_features(prs)
    slide_08_screenshots(prs)
    slide_09_results(prs)
    slide_10_scalability(prs)
    slide_11_value(prs)
    slide_12_closing(prs)

    try:
        prs.save(str(OUT_PATH))
        return OUT_PATH
    except PermissionError:
        # File is likely open in PowerPoint — save to a fresh name instead.
        alt = OUT_PATH.with_name(f"project_presentation_{int(time.time())}.pptx")
        prs.save(str(alt))
        print(f"NOTE: '{OUT_PATH.name}' was locked (open in PowerPoint?). Saved to '{alt.name}'.")
        return alt


if __name__ == "__main__":
    out = build()
    print(f"Presentation generated: {out}")
    print(f"Slides: 12  -  Size: {out.stat().st_size / 1024:.0f} KB")
